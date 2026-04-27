import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_image_safe(slide, img_path, left, top, width, height=None):
    if Path(img_path).exists():
        if height:
            slide.shapes.add_picture(str(img_path), Inches(left), Inches(top),
                                     Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(str(img_path), Inches(left), Inches(top),
                                     Inches(width))
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height or 3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2C, 0x2C, 0x44)
        shape.line.color.rgb = ACCENT_BLUE
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Figure: {Path(img_path).stem}]"
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    fig_dir = PROJECT_ROOT / "figures"

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    add_text_box(slide, 1.5, 1.5, 10, 1.5,
                 "Cross-Site Scripting (XSS) Attack Detection",
                 font_size=36, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 3.2, 10, 0.8,
                 "Using Machine Learning and Hybrid Feature Selection",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 4.5, 10, 0.6,
                 "Team Blackhat",
                 font_size=24, bold=True, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 5.2, 10, 0.5,
                 "Mohammad Ibrahim · Omar Adel · Mohamed Ahmed",
                 font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 6.0, 10, 0.5,
                 "Level 4 Primers · Dr. Wael El Sersy",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 6.6, 10, 0.4,
                 "University of East London (UEL)",
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Slide 2: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 5, 0.7, "Agenda", font_size=30,
                 bold=True, color=ACCENT_BLUE)
    agenda_items = [
        "1. Introduction & Background",
        "2. Problem Statement",
        "3. Literature Review",
        "4. Dataset Overview",
        "5. System Architecture",
        "6. Data Preprocessing",
        "7. Feature Engineering & Analysis",
        "8. Model Selection",
        "9. Results & Performance",
        "10. Live Demo",
        "11. Challenges & Solutions",
        "12. Conclusion & Future Work"
    ]
    add_bullet_list(slide, 1.5, 1.4, 10, 5.5, agenda_items, font_size=18)

    # Slide 3: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 8, 0.7, "What is Cross-Site Scripting (XSS)?",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Major injection vulnerability (OWASP Top Ten)",
        "Attacker injects malicious scripts into trusted websites",
        "Consequences:",
        "  - Session Hijacking",
        "  - Data Exfiltration",
        "  - Phishing & Defacement",
        "Types of XSS:",
        "  - Stored XSS",
        "  - Reflected XSS",
        "  - DOM-based XSS",
        "Need for automated, robust detection mechanisms"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=17)

    # Slide 4: Why ML?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Why Machine Learning?",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Limitations of Rule-Based Systems (WAFs):",
        "  - High false positives",
        "  - Easy to bypass via obfuscation",
        "  - Heavy maintenance burden",
        "Advantages of Machine Learning:",
        "  - Learns complex non-linear patterns",
        "  - Detects zero-day/novel attacks",
        "  - Low maintenance once trained",
        "  - High scalability for real-time traffic"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=17)

    # Slide 5: Literature Review
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Related Work & Literature",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Likarish et al. (2009) - SVM/NB on HTTP parameters -> 92% acc.",
        "Wang et al. (2018) - CNN on character sequences -> 99.2% acc.",
        "Mokbal et al. (2019) - MLP with 42 dynamic features -> 99.0% acc.",
        "Mokbal et al. (2021) - XGBoost + IG-SBS feature selection -> 99.4% acc.",
        "Fang et al. (2018) - LSTM for sequential XSS patterns -> 98.2% acc.",
        "Zhou & Wang (2019) - RF ensemble with HTML/JS/URL -> 98.8% acc.",
        "Our approach: Comparison of RF, XGBoost, and LR on 66 features"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=16)

    # Slide 6: Dataset
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Dataset Overview",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Source: fawaz2015/XSS-dataset (Mokbal et al., 2021)",
        "Total Samples: 138,567",
        "  - Benign: 99,999 (72.1%)",
        "  - Malicious (XSS): 38,568 (27.9%)",
        "Feature Count: 66 pre-selected via IG-SBS",
        "Feature Types:",
        "  - URL structural features",
        "  - HTML DOM features",
        "  - JavaScript behaviour indicators"
    ]
    add_bullet_list(slide, 1.0, 1.3, 6.5, 5.5, items, font_size=16)
    add_image_safe(slide, fig_dir / "01_class_distribution.png", 8.0, 1.5, 4.8, 4.0)

    # Slide 7: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "System Architecture",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Modular ML Pipeline:",
        "1. Data Collection & Preprocessing",
        "2. Hybrid Feature Scaling (StandardScaler)",
        "3. Train/Val/Test Stratified Splitting (70/15/15)",
        "4. Multi-Model Training (RF, XGB, LR)",
        "5. 5-Fold Stratified Cross-Validation",
        "6. Performance Evaluation & Figure Generation",
        "7. Interactive Streamlit Deployment"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=16)

    # Slide 8: Preprocessing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Data Preprocessing",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Duplicate removal (maintaining data integrity)",
        "Missing value imputation (using column medians)",
        "Label encoding (Normal=0, XSS=1)",
        "Feature scaling (fitted on training data only)",
        "Stratified splitting to handle class distribution",
        "Pipeline automation for full reproducibility"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=17)

    # Slide 9: Feature Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Feature Analysis",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    add_image_safe(slide, fig_dir / "02_correlation_heatmap.png", 0.5, 1.3, 6.0, 5.5)
    add_image_safe(slide, fig_dir / "03_feature_importance.png", 6.8, 1.3, 6.0, 5.5)

    # Slide 10: Models
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Model Selection",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Random Forest (RF):",
        "  - Handles non-linear features well; robust to overfitting",
        "Extreme Gradient Boosting (XGBoost):",
        "  - Sequential ensemble; state-of-the-art for tabular data",
        "Logistic Regression (LR):",
        "  - Linear baseline; quantifies non-linear improvement",
        "Evaluation Metrics:",
        "  - Accuracy, Precision, Recall, F1-Score, ROC-AUC"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=16)

    # Slide 11: Results - CM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Results — Confusion Matrices",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    add_image_safe(slide, fig_dir / "04_confusion_matrix_random_forest.png", 0.3, 1.3, 4.1, 3.8)
    add_image_safe(slide, fig_dir / "05_confusion_matrix_xgboost.png", 4.6, 1.3, 4.1, 3.8)
    add_image_safe(slide, fig_dir / "06_confusion_matrix_logistic_regression.png", 8.9, 1.3, 4.1, 3.8)

    # Slide 12: Results - Curves
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Results — ROC & PR Curves",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    add_image_safe(slide, fig_dir / "06_roc_curves_comparison.png", 0.5, 1.3, 6.0, 5.5)
    add_image_safe(slide, fig_dir / "07_precision_recall_curves.png", 6.8, 1.3, 6.0, 5.5)

    # Slide 13: Results - Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Results — Model Comparison",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    add_image_safe(slide, fig_dir / "08_model_comparison.png", 0.5, 1.3, 6.0, 5.5)
    add_image_safe(slide, fig_dir / "09_cv_boxplots.png", 6.8, 1.3, 6.0, 5.5)

    # Slide 14: Results - Groups
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Feature Group Importance",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "URL Features (59.6%): Primary attack vector",
        "JavaScript Features (20.0%): Behavioural analysis",
        "HTML Features (18.1%): Structural analysis",
        "Top features: url_special_characters, js_method_alert, html_tag_script"
    ]
    add_bullet_list(slide, 0.8, 1.3, 5.5, 4, items, font_size=16)
    add_image_safe(slide, fig_dir / "10_feature_group_analysis.png", 7.0, 1.3, 5.8, 4.5)

    # Slide 15: Demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Live Demo",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Built with Streamlit for rapid deployment",
        "Features:",
        "  - Real-time prediction from manual input",
        "  - Interactive model selection",
        "  - Batch testing on random test samples",
        "  - Visual performance monitoring",
        "Provides an easy-to-use interface for security analysts"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=17)

    # Slide 16: Challenges
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Challenges & Solutions",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "Challenge: Large NVIDIA dependencies in XGBoost",
        "Solution: Docker cleanup scripts to remove CUDA overhead",
        "Challenge: CSV structure changes (Label vs class)",
        "Solution: Robust automated header detection in preprocessing",
        "Challenge: Reproducibility across environments",
        "Solution: Full Dockerisation and fixed random seeds"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=16)

    # Slide 17: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.8, 0.4, 10, 0.7, "Conclusion & Future Work",
                 font_size=30, bold=True, color=ACCENT_BLUE)
    items = [
        "ML successfully detects XSS with >99% accuracy",
        "Tree-based ensembles (XGB, RF) significantly outperform LR",
        "Feature engineering is key to performance",
        "Future Work:",
        "  - Deep learning (LSTMs) for raw payload analysis",
        "  - Adversarial robustness testing",
        "  - Real-time browser plugin integration"
    ]
    add_bullet_list(slide, 1.0, 1.3, 11, 5.5, items, font_size=16)

    # Slide 18: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 1.5, 2.5, 10, 1.5, "Thank You!",
                 font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 4.2, 10, 0.6, "Questions?",
                 font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 5.3, 10, 0.5, "Team Blackhat",
                 font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.5, 5.9, 10, 0.5, "University of East London",
                 font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    output_path = PROJECT_ROOT / "presentation" / "Blackhat_Presentation.pptx"
    prs.save(str(output_path))
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
